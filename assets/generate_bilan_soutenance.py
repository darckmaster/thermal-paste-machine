"""
Génère les planches de BILAN et de PORTAGE CNC pour la soutenance blanche.

Complète `presentation.pptx` (généré par generate_presentation.py) : ce fichier-ci ne
produit que les 7 planches nouvelles, à copier dans le deck existant. Style et helpers
repris tels quels de generate_presentation.py — la charte reste la même.

Exécuter depuis le dossier assets/ :  python generate_bilan_soutenance.py
Sortie : bilan_soutenance.pptx

⚠️ Les chiffres de ce fichier sont MESURÉS sur le dépôt au 2026-08-04, pas estimés.
Pour les rafraîchir avant une autre soutenance, relancer les commandes indiquées en
commentaire au-dessus de chaque constante.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from generate_presentation import (
    W, H,
    C_WHITE, C_BLACK, C_ACCENT, C_ACCENT2, C_GRAY, C_LGRAY,
    C_GREEN, C_ORANGE, C_RED,
    FONT,
    add_rect, add_text_box, add_para, set_slide_bg, header_bar, content_box,
)

TOTAL_SLIDES = 9

# --------------------------------------------------------------------------- chiffres
# Mesurés le 2026-08-04. Commandes de rafraîchissement :
#   cat modules/*.py gui/*.py main.py | grep -c ''      -> lignes de code
#   cat tests/test_*.py | grep -c ''                    -> lignes de test
#   grep -h -c '^def test_' tests/test_*.py             -> fonctions de test
#   git rev-list --count HEAD                           -> commits
#   grep -o '_local_cfg.get("[a-z_]*"' modules/config.py | sort -u | wc -l
LIGNES_CODE = 9971
LIGNES_TESTS = 5754
NB_TESTS = 300
NB_COMMITS = 75
NB_VERSIONS = 12
NB_PARAMS_EXTERNALISES = 25
LIGNES_A_AJUSTER = 9


def _slide(prs, titre, sous_titre=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    header_bar(slide, titre, sous_titre)
    return slide


def _bloc(slide, x, y, largeur, titre, lignes, couleur=C_ACCENT2, hauteur=Cm(7.0)):
    """Un bloc de contenu : filet coloré à gauche, titre, puis des lignes."""
    add_rect(slide, x, y, Pt(4), hauteur, fill_rgb=couleur)
    boite = slide.shapes.add_textbox(x + Pt(10), y, largeur - Pt(10), hauteur)
    tf = boite.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = titre
    run.font.name = FONT
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = couleur

    for texte, gras, coul, taille in lignes:
        add_para(tf, texte, font_size=Pt(taille), bold=gras, color=coul,
                 space_before=Pt(5))
    return tf


def _chiffre(slide, x, y, valeur, libelle, couleur=C_ACCENT):
    """Un chiffre-clé mis en avant."""
    add_text_box(slide, valeur, x, y, Cm(6.0), Cm(1.5),
                 font_size=Pt(34), bold=True, color=couleur, align=PP_ALIGN.CENTER)
    add_text_box(slide, libelle, x, y + Cm(1.5), Cm(6.0), Cm(1.2),
                 font_size=Pt(11), color=C_GRAY, align=PP_ALIGN.CENTER)


def _pied(slide, texte, couleur=C_ACCENT):
    """Bandeau de conclusion en bas de planche — la phrase à retenir."""
    add_rect(slide, Cm(1.2), Cm(16.9), Cm(31.0), Cm(1.5), fill_rgb=C_LGRAY)
    add_rect(slide, Cm(1.2), Cm(16.9), Pt(5), Cm(1.5), fill_rgb=couleur)
    add_text_box(slide, texte, Cm(1.6), Cm(17.15), Cm(30.2), Cm(1.0),
                 font_size=Pt(14), bold=True, color=C_BLACK)


# ===========================================================================
# SYNOPTIQUE DU PROCESSUS MÉTIER
# ===========================================================================

def _etape(slide, x, y, largeur, hauteur, numero, titre, detail, couleur):
    """Une étape du synoptique : pastille numérotée, titre, explication.

    Le numéro n'est pas décoratif : il donne au jury un point d'accroche pour poser une
    question sur une étape précise sans avoir à la décrire.
    """
    from pptx.enum.shapes import MSO_SHAPE

    boite = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, largeur, hauteur)
    boite.fill.solid()
    boite.fill.fore_color.rgb = C_WHITE
    boite.line.color.rgb = couleur
    boite.line.width = Pt(1.75)
    boite.shadow.inherit = False
    boite.text_frame.text = ""

    # Pastille du numéro, en haut à gauche du cadre
    pastille = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Cm(0.25), y + Cm(0.22), Cm(0.85), Cm(0.85)
    )
    pastille.fill.solid()
    pastille.fill.fore_color.rgb = couleur
    pastille.line.fill.background()
    pastille.shadow.inherit = False
    tf_num = pastille.text_frame
    tf_num.word_wrap = False
    p_num = tf_num.paragraphs[0]
    p_num.alignment = PP_ALIGN.CENTER
    r_num = p_num.add_run()
    r_num.text = str(numero)
    r_num.font.name = FONT
    r_num.font.size = Pt(12)
    r_num.font.bold = True
    r_num.font.color.rgb = C_WHITE

    add_text_box(slide, titre, x + Cm(1.2), y + Cm(0.18), largeur - Cm(1.4), Cm(1.1),
                 font_size=Pt(12), bold=True, color=couleur)

    boite_detail = slide.shapes.add_textbox(
        x + Cm(0.3), y + Cm(1.25), largeur - Cm(0.6), hauteur - Cm(1.4)
    )
    tf = boite_detail.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = detail
    run.font.name = FONT
    run.font.size = Pt(10)
    run.font.color.rgb = C_BLACK
    return boite


def _fleche(slide, x, y, largeur, hauteur, couleur, vers_le_bas=False):
    from pptx.enum.shapes import MSO_SHAPE

    forme = MSO_SHAPE.DOWN_ARROW if vers_le_bas else MSO_SHAPE.RIGHT_ARROW
    f = slide.shapes.add_shape(forme, x, y, largeur, hauteur)
    f.fill.solid()
    f.fill.fore_color.rgb = couleur
    f.line.fill.background()
    f.shadow.inherit = False
    return f


def slide_synoptique(prs):
    """Le processus métier, de bout en bout, sur une seule planche."""
    from pptx.enum.shapes import MSO_SHAPE

    slide = _slide(
        prs, "Le processus métier couvert par le logiciel",
        "On prépare une référence UNE fois — on la produit ensuite autant de fois qu'on veut"
    )

    # ---------------------------------------------------------------- bandeau 1
    add_text_box(slide, "PRÉPARATION  ·  une seule fois par référence produit",
                 Cm(1.5), Cm(3.25), Cm(20), Cm(0.7),
                 font_size=Pt(13), bold=True, color=C_ACCENT2)

    y1, h1, w1, a1 = Cm(4.0), Cm(3.1), Cm(6.9), Cm(1.05)
    x = Cm(1.5)
    etapes_prep = [
        ("Nommer le produit",
         "L'opérateur saisit la référence, ou la reprend dans la liste des produits déjà connus."),
        ("Photographier le plateau",
         "Le logiciel reconnaît seul les emplacements et signale ceux qui sont mal montés."),
        ("Tracer les cordons",
         "Au doigt, sur UN seul emplacement agrandi à l'écran."),
        ("Enregistrer",
         "Le tracé est reporté automatiquement sur tous les autres emplacements."),
    ]
    for i, (titre, detail) in enumerate(etapes_prep, 1):
        _etape(slide, x, y1, w1, h1, i, titre, detail, C_ACCENT2)
        x += w1
        if i < len(etapes_prep):
            _fleche(slide, x + Cm(0.15), y1 + Cm(1.1), a1 - Cm(0.3), Cm(0.9), C_LGRAY)
            x += a1

    # ---------------------------------------------------------------- le fichier
    _fleche(slide, Cm(16.2), Cm(7.25), Cm(0.9), Cm(0.75), C_ORANGE, vers_le_bas=True)

    fichier = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Cm(7.0), Cm(8.05), Cm(19.5), Cm(1.35))
    fichier.fill.solid()
    fichier.fill.fore_color.rgb = C_ORANGE
    fichier.line.fill.background()
    fichier.shadow.inherit = False
    tf_f = fichier.text_frame
    tf_f.word_wrap = True
    p_f = tf_f.paragraphs[0]
    p_f.alignment = PP_ALIGN.CENTER
    r_f = p_f.add_run()
    r_f.text = ("FICHIER PLATEAU  —  produit, emplacements, cordons, réglages   "
                "·   rejouable sans rien retracer")
    r_f.font.name = FONT
    r_f.font.size = Pt(12)
    r_f.font.bold = True
    r_f.font.color.rgb = C_WHITE

    _fleche(slide, Cm(16.2), Cm(9.6), Cm(0.9), Cm(0.75), C_GREEN, vers_le_bas=True)

    # ---------------------------------------------------------------- bandeau 2
    add_text_box(slide, "PRODUCTION  ·  à chaque plateau, en série",
                 Cm(1.5), Cm(10.55), Cm(20), Cm(0.7),
                 font_size=Pt(13), bold=True, color=C_GREEN)

    y2, h2, w2, a2 = Cm(11.3), Cm(3.1), Cm(5.5), Cm(0.8)
    x = Cm(1.5)
    etapes_prod = [
        ("Lancer",
         "Un bouton. La machine se met d'aplomb et en position d'observation."),
        ("Choisir le plateau",
         "Dans la liste des références déjà préparées."),
        ("Désigner les pièces",
         "Le logiciel montre les emplacements ; l'opérateur touche ceux qui sont garnis."),
        ("Déposer",
         "Dépose automatique, avec avancement, pause et arrêt d'urgence."),
        ("Rapport",
         "Photo de fin et PDF de traçabilité : ce qui a été déposé, et où."),
    ]
    for i, (titre, detail) in enumerate(etapes_prod, 1):
        _etape(slide, x, y2, w2, h2, i, titre, detail, C_GREEN)
        x += w2
        if i < len(etapes_prod):
            _fleche(slide, x + Cm(0.1), y2 + Cm(1.1), a2 - Cm(0.2), Cm(0.9), C_LGRAY)
            x += a2

    # ---------------------------------------------------------------- boucle de série
    # Flèche de RETOUR sous la bande, plutôt qu'un bloc à droite : un bloc à cet endroit
    # chevaucherait la dernière étape, et la boucle est de toute façon plus lisible
    # dessinée comme un retour au début qu'énoncée dans un cadre.
    retour = slide.shapes.add_shape(
        MSO_SHAPE.LEFT_ARROW, Cm(1.5), Cm(14.85), Cm(30.7), Cm(1.15)
    )
    retour.fill.solid()
    retour.fill.fore_color.rgb = C_LGRAY
    retour.line.fill.background()
    retour.shadow.inherit = False
    tf_r = retour.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.alignment = PP_ALIGN.CENTER
    r_r = p_r.add_run()
    r_r.text = ("PLATEAU SUIVANT  —  on repart à l'étape 1, sans rien reconfigurer "
                "ni retracer")
    r_r.font.name = FONT
    r_r.font.size = Pt(12)
    r_r.font.bold = True
    r_r.font.color.rgb = C_GREEN

    _pied(slide,
          "Le tracé est fait une fois pour une référence ; ensuite, produire un plateau "
          "ne demande que trois gestes au doigt.", C_GREEN)
    return slide


# ===========================================================================
# ARCHITECTURE LOGICIELLE — la carte des modules
# ===========================================================================

# Nombre de lignes relevé le 2026-08-04 :
#   for f in main.py modules/*.py gui/*.py; do echo "$f $(grep -c '' $f)"; done
#
# Le compte de lignes n'est pas de la vanité : il montre d'un coup d'œil où est
# concentré l'effort — la vision et le tracé pèsent à eux deux un quart du logiciel.
MODULES_NOYAU = [
    ("config.py", 254, "Tous les paramètres, dont ceux propres à chaque machine"),
    ("camera.py", 181, "Capture USB, choix du périphérique, image fraîche garantie"),
    ("vision.py", 1225, "ArUco, repère du plateau, reconnaissance des emplacements"),
    ("calibration.py", 398, "Correction de la distorsion de l'objectif (mire ChArUco)"),
    ("machine.py", 293, "Dialogue G-code avec Marlin, par liaison série"),
    ("path_planner.py", 606, "Trajectoire : ordre des zones, cordons, contrôle de course"),
    ("preparation.py", 767, "Modèle d'un plateau et son enregistrement (JSON versionné)"),
    ("reporter.py", 383, "Rapport PDF de traçabilité (fpdf2)"),
]

MODULES_IHM = [
    ("main.py", 36, "Point d'entrée de l'application"),
    ("app.py", 481, "Fenêtre principale et navigation entre les écrans"),
    ("workers.py", 127, "Mise en position machine, hors du fil d'affichage"),
    ("dialogs.py", 673, "Boîtes de dialogue partagées (produit, paramètres, dépose)"),
    ("screen_capture.py", 589, "Accueil : aperçu caméra, choix du matériel, homing"),
    ("screen_plateau.py", 619, "Créer un plateau : photo et diagnostic des emplacements"),
    ("screen_cordons.py", 948, "Tracé des cordons au doigt, report sur tous les emplacements"),
    ("screen_execution.py", 950, "Cycle de dépose multi-zones, de bout en bout"),
    ("screen_calibration.py", 489, "Calibration optique : capture des poses, calcul"),
]

MODULES_HISTORIQUE = [
    ("screen_zone.py", 476, "Ancien cycle mono-zone — retrait planifié"),
    ("screen_run.py", 320, "Ancienne exécution — retrait planifié"),
    ("screen_report.py", 155, "Ancien rapport — retrait planifié"),
]


# Géométrie des cartes. Le pas est contraint par la colonne la plus chargée (9 cartes)
# et par le bandeau de conclusion, qui commence à 16,9 cm : 3,95 + 8 × 1,39 + 1,25 =
# 16,32 cm, soit un peu plus d'un demi-centimètre de marge. Toute carte plus haute
# déborderait sur le bandeau — c'est ce qui est arrivé à la première version.
_CARTE_H = Cm(1.25)
_CARTE_ECART = Cm(0.14)
_CARTE_L = Cm(10.1)
_COLONNES = (Cm(1.2), Cm(11.85), Cm(22.5))
_CARTE_Y0 = Cm(3.95)


def _carte_module(slide, x, y, nom, lignes, description, couleur, estompe=False):
    """Une carte de module : filet coloré à gauche, nom, volume, rôle en une ligne."""
    add_rect(slide, x, y, _CARTE_L, _CARTE_H,
             fill_rgb=C_WHITE, line_rgb=C_LGRAY, line_width=Pt(0.75))
    add_rect(slide, x, y, Cm(0.2), _CARTE_H, fill_rgb=couleur)

    add_text_box(slide, nom, x + Cm(0.42), y + Cm(0.07), _CARTE_L - Cm(2.9), Cm(0.6),
                 font_size=Pt(11), bold=True, color=C_GRAY if estompe else couleur)
    add_text_box(slide, f"{lignes} l.", x + _CARTE_L - Cm(2.2), y + Cm(0.1),
                 Cm(1.8), Cm(0.55), font_size=Pt(8.5), color=C_GRAY,
                 align=PP_ALIGN.RIGHT)
    add_text_box(slide, description, x + Cm(0.42), y + Cm(0.6),
                 _CARTE_L - Cm(0.65), Cm(0.6),
                 font_size=Pt(8.5), color=C_GRAY if estompe else C_BLACK)


def _colonne_cartes(slide, colonne, y, entrees, couleur, accent_sur=None,
                    estompe=False):
    """Empile une liste de cartes dans une colonne, et rend le y atteint."""
    for nom, n, desc in entrees:
        c = C_ACCENT if nom == accent_sur else couleur
        _carte_module(slide, _COLONNES[colonne], y, nom, n, desc, c, estompe)
        y += _CARTE_H + _CARTE_ECART
    return y


def slide_modules(prs):
    """La carte complète des modules du logiciel."""
    total = len(MODULES_NOYAU) + len(MODULES_IHM) + len(MODULES_HISTORIQUE)
    slide = _slide(prs, f"Architecture logicielle — les {total} modules",
                   "Une responsabilité par module : le cœur métier ne connaît "
                   "ni l'écran ni la machine")

    add_text_box(slide, "CŒUR MÉTIER  ·  modules/", _COLONNES[0], Cm(3.3),
                 Cm(10.0), Cm(0.6), font_size=Pt(12), bold=True, color=C_ACCENT2)
    _colonne_cartes(slide, 0, _CARTE_Y0, MODULES_NOYAU, C_ACCENT2,
                    accent_sur="config.py")

    add_text_box(slide, "INTERFACE ET PILOTAGE  ·  gui/", _COLONNES[1], Cm(3.3),
                 Cm(10.0), Cm(0.6), font_size=Pt(12), bold=True, color=C_GREEN)
    _colonne_cartes(slide, 1, _CARTE_Y0, MODULES_IHM, C_GREEN, accent_sur="main.py")

    # Colonne 3 : le cycle historique, volontairement estompé. Il est encore là — c'est
    # le seul chemin déjà validé jusqu'à la dépose réelle — mais son retrait est planifié.
    add_text_box(slide, "CYCLE HISTORIQUE  ·  retrait planifié", _COLONNES[2], Cm(3.3),
                 Cm(10.0), Cm(0.6), font_size=Pt(12), bold=True, color=C_GRAY)
    y = _colonne_cartes(slide, 2, _CARTE_Y0, MODULES_HISTORIQUE, C_LGRAY, estompe=True)

    add_text_box(slide,
                 "Remplacé par le nouveau cycle multi-zones, mais conservé tant que "
                 "celui-ci n'a pas été validé sur la machine : on ne retire pas ce qui "
                 "fonctionne avant d'avoir prouvé son remplaçant.",
                 _COLONNES[2] + Cm(0.42), y + Cm(0.3), _CARTE_L - Cm(0.8), Cm(2.4),
                 font_size=Pt(10), italic=True, color=C_GRAY)

    _pied(slide,
          "Un module par responsabilité : la vision ne pilote rien, la machine ne voit "
          "rien, et l'interface ne calcule rien.")
    return slide


# ===========================================================================
# PARTIE 1 — BILAN
# ===========================================================================

def slide_acquis(prs):
    """Ce qui fonctionne aujourd'hui, de bout en bout."""
    slide = _slide(prs, "Bilan — ce qui fonctionne aujourd'hui",
                   "Chaîne complète opérationnelle sur la machine de prototypage")

    n = (False, C_BLACK, 13)
    _bloc(slide, Cm(1.2), Cm(3.6), Cm(15.0), "VISION ET REPÉRAGE", [
        ("Détection ArUco des 4 coins du plateau", *n),
        ("Repère orthonormé, origine sur un marqueur", *n),
        ("Fonctionne avec 2 coins sur 4 seulement", *n),
        ("   → mode nominal, la caméra ne cadre pas tout", False, C_GRAY, 11),
        ("Reconnaissance automatique des emplacements", *n),
        ("Diagnostic des défauts de montage", *n),
        ("Calibration optique ChArUco intégrée", *n),
    ], C_ACCENT2)

    _bloc(slide, Cm(17.4), Cm(3.6), Cm(15.0), "PRÉPARATION ET EXÉCUTION", [
        ("Tracé des cordons à l'écran tactile", *n),
        ("Un tracé rejoué sur tous les emplacements", *n),
        ("Enregistrement, rechargement, reprise", *n),
        ("Cycle de dépose complet, guidé, multi-zones", *n),
        ("Pilotage machine : homing, déplacements, arrêt", *n),
        ("Suivi temps réel, pause, arrêt d'urgence", *n),
        ("Photo de fin et rapport PDF de traçabilité", *n),
    ], C_GREEN)

    _bloc(slide, Cm(1.2), Cm(11.2), Cm(31.2), "CE QUE L'OPÉRATEUR FAIT, CONCRÈTEMENT", [
        ("Un bouton → la machine se repère seule, il désigne les pièces présentes, "
         "confirme, et la dépose s'exécute sous ses yeux avec un rapport à la clé.",
         False, C_BLACK, 14),
        ("Aucune ligne de code, aucun réglage de fichier : tout se fait au doigt sur "
         "l'écran 7 pouces.", False, C_GRAY, 12),
    ], C_ACCENT, hauteur=Cm(4.8))

    _pied(slide, "Le processus métier complet est couvert, du plateau nu au rapport signé.",
          C_GREEN)
    return slide


def slide_chiffres(prs):
    """Le chiffrage du travail."""
    slide = _slide(prs, "Bilan — le chiffrage",
                   "Mesuré sur le dépôt, au 4 août")

    y = Cm(4.2)
    _chiffre(slide, Cm(1.5), y, f"{LIGNES_CODE:,}".replace(",", " "),
             "lignes de code applicatif")
    _chiffre(slide, Cm(8.0), y, f"{LIGNES_TESTS:,}".replace(",", " "),
             "lignes de tests automatiques", C_GREEN)
    _chiffre(slide, Cm(14.5), y, str(NB_TESTS),
             "tests, exécutés en 35 s", C_GREEN)
    _chiffre(slide, Cm(21.0), y, str(NB_COMMITS),
             "commits documentés")
    _chiffre(slide, Cm(27.5), y, str(NB_VERSIONS),
             "versions livrées")

    n = (False, C_BLACK, 13)
    _bloc(slide, Cm(1.2), Cm(8.4), Cm(15.0), "UN RAPPORT TESTS / CODE DE 0,58", [
        ("Plus d'une ligne de test pour deux lignes de code.", *n),
        ("Sur la dépose — la fonction critique — les tests sont", *n),
        ("validés par MUTATION : on casse volontairement le", *n),
        ("code pour vérifier qu'ils réagissent.", *n),
        ("→ 3 tests décoratifs démasqués alors que tout", False, C_ACCENT, 12),
        ("   était au vert.", False, C_ACCENT, 12),
    ], C_GREEN, hauteur=Cm(6.2))

    _bloc(slide, Cm(17.4), Cm(8.4), Cm(15.0), "UNE TRAÇABILITÉ COMPLÈTE", [
        ("Chaque décision technique est consignée avec son", *n),
        ("MOTIF — ce qui empêche de la « corriger » plus tard", *n),
        ("en croyant réparer un oubli.", *n),
        ("2 manuels (opérateur, maintenance), un dossier de", *n),
        ("conception, un journal de session.", *n),
        ("→ le projet est reprenable par un tiers.", False, C_ACCENT, 12),
    ], C_ACCENT2, hauteur=Cm(6.2))

    _pied(slide, "Trois mois de travail, une base de code testée et documentée — "
                 "pas un prototype jetable.")
    return slide


def slide_reste_a_faire(prs):
    """Ce qui reste — honnête."""
    slide = _slide(prs, "Bilan — ce qui reste à faire",
                   "Vue honnête : trois chantiers, dont un seul est bloquant")

    n = (False, C_BLACK, 13)
    _bloc(slide, Cm(1.2), Cm(3.6), Cm(10.0), "1. VALIDER SUR MACHINE", [
        ("Le cycle multi-zones complet n'a encore", *n),
        ("jamais tourné sur la machine : il a été", *n),
        ("écrit et testé à la maison.", *n),
        ("", *n),
        ("Coût : une demi-journée d'atelier.", False, C_ACCENT, 13),
        ("Sans risque : mode « à blanc », la buse", False, C_GRAY, 12),
        ("ne descend pas et n'extrude rien.", False, C_GRAY, 12),
    ], C_RED, hauteur=Cm(7.6))

    _bloc(slide, Cm(12.0), Cm(3.6), Cm(10.0), "2. RÉGLER L'EXTRUSION", [
        ("La pâte est très visqueuse : elle sort en", *n),
        ("retard au démarrage et continue après", *n),
        ("l'arrêt.", *n),
        ("", *n),
        ("Les deux compensations sont CODÉES et", False, C_ACCENT, 13),
        ("paramétrables ; il reste à les régler à", False, C_ACCENT, 13),
        ("l'œil, avec la pâte réelle.", False, C_ACCENT, 13),
    ], C_ORANGE, hauteur=Cm(7.6))

    _bloc(slide, Cm(22.8), Cm(3.6), Cm(9.6), "3. FINITIONS", [
        ("Bascule du nouveau cycle en point", *n),
        ("d'entrée unique, après validation.", *n),
        ("Retrait de l'ancien cycle mono-zone.", *n),
        ("", *n),
        ("3 points de robustesse identifiés,", False, C_GRAY, 12),
        ("audités, chiffrés.", False, C_GRAY, 12),
        ("Calibration optique sur le matériel réel.", False, C_GRAY, 12),
    ], C_ACCENT2, hauteur=Cm(7.6))

    _bloc(slide, Cm(1.2), Cm(11.8), Cm(31.2), "POURQUOI CE N'EST PAS INQUIÉTANT", [
        ("Aucun de ces chantiers ne demande d'écrire une fonctionnalité nouvelle : "
         "ce sont des réglages, des mesures et des retraits.", False, C_BLACK, 14),
        ("Les 11 mesures machine restantes sont NOMMÉES, NUMÉROTÉES et rappelées au "
         "début de chaque séance de travail — aucune ne peut se perdre.",
         False, C_BLACK, 13),
    ], C_ACCENT, hauteur=Cm(4.4))

    _pied(slide, "Ce qui reste tient en réglages et en mesures — plus en développement.",
          C_ORANGE)
    return slide


# ===========================================================================
# PARTIE 2 — PORTAGE CNC
# ===========================================================================

def slide_forge(prs):
    """Le PoC comme forge de la machine finale."""
    slide = _slide(prs, "Portage CNC — le prototype était la forge",
                   "Tout ce qui a été construit sur la machine d'essai est réutilisable tel quel")

    n = (False, C_BLACK, 13)
    ok = (False, C_GREEN, 13)

    _bloc(slide, Cm(1.2), Cm(3.6), Cm(15.4), "IDENTIQUE SUR LES DEUX MACHINES", [
        ("Firmware Marlin — même dialecte G-code", *ok),
        ("Caméra Philips SPC 1330NC — même modèle", *ok),
        ("Raspberry Pi + écran tactile 7 pouces", *ok),
        ("Marqueurs ArUco, même dictionnaire", *ok),
        ("Toute la vision, toute l'interface, toute la", *ok),
        ("logique de trajectoire", *ok),
        ("", *n),
        ("Ce choix a été fait en juillet, volontairement :", False, C_GRAY, 11),
        ("une seule base de code pour les deux machines.", False, C_GRAY, 11),
    ], C_GREEN, hauteur=Cm(8.2))

    _bloc(slide, Cm(17.8), Cm(3.6), Cm(14.6), "CE QUI DIFFÈRE", [
        ("La caméra est fixée sur la seringue au lieu", *n),
        ("du bâti.", *n),
        ("", *n),
        ("Conséquence unique : la position depuis", False, C_ORANGE, 13),
        ("laquelle on photographie devient une", False, C_ORANGE, 13),
        ("inconnue à mesurer.", False, C_ORANGE, 13),
        ("", *n),
        ("Elle est DÉJÀ un paramètre du logiciel,", False, C_ACCENT, 12),
        ("prévu et écrit depuis le 3 août.", False, C_ACCENT, 12),
    ], C_ORANGE, hauteur=Cm(8.2))

    _bloc(slide, Cm(1.2), Cm(12.4), Cm(31.2), "LE POINT DE CONCEPTION QUI REND LE PORTAGE POSSIBLE", [
        ("Le repère du plateau est ancré sur les MARQUEURS, jamais sur la mécanique. "
         "La vision est donc indépendante du châssis, de la course et de la cinématique.",
         False, C_BLACK, 14),
    ], C_ACCENT, hauteur=Cm(3.6))

    _pied(slide, "La machine d'essai n'était pas une maquette : c'est la même machine, "
                 "à la mécanique près.", C_GREEN)
    return slide


def slide_effort(prs):
    """L'effort de portage, chiffré."""
    slide = _slide(prs, "Portage CNC — l'effort, chiffré",
                   "Combien de code faut-il réellement toucher ?")

    y = Cm(4.0)
    _chiffre(slide, Cm(2.0), y, "0", "ligne de logique à réécrire", C_GREEN)
    _chiffre(slide, Cm(10.0), y, f"{LIGNES_A_AJUSTER}", "valeurs numériques à ajuster", C_ORANGE)
    _chiffre(slide, Cm(18.0), y, "0,09 %", "du code applicatif", C_ORANGE)
    _chiffre(slide, Cm(26.0), y, f"{NB_PARAMS_EXTERNALISES}",
             "réglages déjà hors du code", C_GREEN)

    n = (False, C_BLACK, 13)
    _bloc(slide, Cm(1.2), Cm(8.2), Cm(15.4), "UN SEUL FICHIER CONCERNÉ", [
        ("Toutes les grandeurs qui dépendent de la machine", *n),
        ("sont regroupées dans un fichier de configuration", *n),
        ("PROPRE À CHAQUE MACHINE, hors du code suivi.", *n),
        ("", *n),
        ("Port série, caméra, taille du plateau, origine,", False, C_GRAY, 12),
        ("course des axes, position de prise de vue,", False, C_GRAY, 12),
        ("vitesses, marges de sécurité.", False, C_GRAY, 12),
        ("", *n),
        ("Aucune valeur machine n'est écrite en dur ailleurs", False, C_ACCENT, 12),
        ("dans le code — vérifié.", False, C_ACCENT, 12),
    ], C_GREEN, hauteur=Cm(8.0))

    _bloc(slide, Cm(17.8), Cm(8.2), Cm(14.6), "LES 9 VALEURS À AJUSTER", [
        ("Position du plateau dans le repère machine (2)", *n),
        ("Hauteurs de dépose et de transit (3)", *n),
        ("Vitesses de déplacement et d'extrusion (3)", *n),
        ("Hauteur caméra (1)", *n),
        ("", *n),
        ("Ce sont des MESURES, pas des décisions de", False, C_ACCENT, 13),
        ("conception. Chacune se relève en quelques", False, C_ACCENT, 13),
        ("minutes, machine sous tension.", False, C_ACCENT, 13),
    ], C_ORANGE, hauteur=Cm(8.0))

    _pied(slide, "Porter le logiciel sur la CNC, c'est remplir un fichier de mesures — "
                 "pas réécrire un programme.", C_GREEN)
    return slide


def slide_manipulations(prs):
    """Les manipulations restantes, avec leur coût."""
    slide = _slide(prs, "Portage CNC — les manipulations restantes",
                   "Une liste à cocher, pas une exploration")

    n = (False, C_BLACK, 12)
    _bloc(slide, Cm(1.2), Cm(3.6), Cm(15.4), "MATÉRIEL — LE VRAI CHANTIER", [
        ("Finir le câblage capteurs et moteurs", *n),
        ("Première mise sous tension", *n),
        ("Réglage : sens moteurs, courants, pas/mm", *n),
        ("", *n),
        ("C'est le seul poste dont la durée n'est pas", False, C_RED, 12),
        ("connue d'avance. Il ne dépend pas du logiciel.", False, C_RED, 12),
        ("", *n),
        ("Mécanique et carte déjà assemblées,", False, C_GREEN, 12),
        ("firmware déjà flashé.", False, C_GREEN, 12),
    ], C_ORANGE, hauteur=Cm(7.8))

    _bloc(slide, Cm(17.8), Cm(3.6), Cm(14.6), "MESURES — 7 RELEVÉS, ~1 H AU TOTAL", [
        ("Taille du plateau (mètre)", *n),
        ("Origine du plateau (pointe au homing)", *n),
        ("Hauteur de la pointe", *n),
        ("Sens réel des axes", *n),
        ("Course utile des axes", *n),
        ("Position de prise de vue", *n),
        ("Calibration optique (15 poses)", *n),
        ("", *n),
        ("Toutes déjà faites une fois sur le prototype :", False, C_ACCENT, 12),
        ("le mode opératoire est écrit dans le manuel.", False, C_ACCENT, 12),
    ], C_ACCENT2, hauteur=Cm(7.8))

    _bloc(slide, Cm(1.2), Cm(12.0), Cm(31.2), "L'ACQUIS LE PLUS PRÉCIEUX DU PROTOTYPE", [
        ("Sur la machine d'essai, chacune de ces grandeurs a coûté une séance de "
         "découverte : il a fallu comprendre qu'elle existait, qu'elle comptait, et "
         "comment la mesurer sans se tromper.", False, C_BLACK, 14),
        ("Sur la CNC, ce n'est plus une exploration : c'est une liste à cocher, avec un "
         "mode opératoire éprouvé et les pièges déjà documentés.", False, C_ACCENT, 14),
    ], C_ACCENT, hauteur=Cm(4.2))

    _pied(slide, "Le prototype a transformé des inconnues en procédure.", C_GREEN)
    return slide


def slide_calendrier(prs):
    """Pourquoi fin août est crédible."""
    slide = _slide(prs, "Portage CNC — pourquoi fin août est tenable",
                   "Le chemin critique est matériel, pas logiciel")

    n = (False, C_BLACK, 13)
    _bloc(slide, Cm(1.2), Cm(3.6), Cm(15.4), "CE QUI EST DÉJÀ DERRIÈRE NOUS", [
        ("Mécanique CNC assemblée", *n),
        ("Carte de commande installée", *n),
        ("Firmware Marlin flashé", *n),
        ("Logiciel complet, testé, documenté", *n),
        ("Procédures de mesure écrites et éprouvées", *n),
        ("Configuration par machine déjà prévue", *n),
        ("", *n),
        ("Le logiciel n'attend plus que la machine.", False, C_GREEN, 14),
    ], C_GREEN, hauteur=Cm(7.4))

    _bloc(slide, Cm(17.8), Cm(3.6), Cm(14.6), "LE RISQUE, NOMMÉ", [
        ("Le seul poste à durée incertaine est la mise", *n),
        ("au point électromécanique de la CNC.", *n),
        ("", *n),
        ("Atténuation : le prototype reste opérationnel", False, C_ACCENT, 13),
        ("et permet de continuer à faire progresser le", False, C_ACCENT, 13),
        ("logiciel en parallèle, sans dépendre de la CNC.", False, C_ACCENT, 13),
        ("", *n),
        ("Une démonstration sur le prototype est déjà", False, C_GRAY, 12),
        ("acceptée pour la soutenance finale.", False, C_GRAY, 12),
    ], C_ORANGE, hauteur=Cm(7.4))

    _bloc(slide, Cm(1.2), Cm(11.6), Cm(31.2), "ENCHAÎNEMENT PROPOSÉ", [
        ("1.  Valider le cycle complet sur le prototype, à blanc  →  une demi-journée",
         False, C_BLACK, 13),
        ("2.  Régler l'extrusion avec la pâte réelle  →  une à deux séances d'essais",
         False, C_BLACK, 13),
        ("3.  Terminer le câblage CNC et la mise sous tension  →  poste à durée variable",
         False, C_BLACK, 13),
        ("4.  Relever les 7 mesures et remplir le fichier de configuration  →  une heure",
         False, C_BLACK, 13),
        ("5.  Rejouer le même cycle sur la CNC  →  aucun développement",
         False, C_ACCENT, 13),
    ], C_ACCENT, hauteur=Cm(4.8))

    _pied(slide, "Deux machines fonctionnelles ne demandent plus qu'un logiciel : "
                 "elles demandent du câblage et sept mesures.", C_GREEN)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_synoptique(prs)
    slide_modules(prs)
    slide_acquis(prs)
    slide_chiffres(prs)
    slide_reste_a_faire(prs)
    slide_forge(prs)
    slide_effort(prs)
    slide_manipulations(prs)
    slide_calendrier(prs)

    # Écrire sous un nom de repli si le fichier est verrouillé — c'est le cas dès qu'il
    # est ouvert dans PowerPoint, ce qui arrive naturellement quand on itère sur le
    # rendu. Échouer là-dessus ferait perdre la génération entière pour rien.
    sortie = "bilan_soutenance.pptx"
    try:
        prs.save(sortie)
    except PermissionError:
        import datetime
        sortie = f"bilan_soutenance_{datetime.datetime.now():%H%M%S}.pptx"
        prs.save(sortie)
        print("(le fichier principal etait ouvert — ecriture sous un nom de repli)")
    print(f"Planches generees : {sortie}  ({TOTAL_SLIDES} slides)")


if __name__ == "__main__":
    main()
