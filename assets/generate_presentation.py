"""
Génère la présentation PowerPoint du projet machine de dépose de pâte thermique.
Style volontairement sobre (fond blanc, texte noir/gris, une couleur d'accent)
pour faciliter les retouches manuelles.
Exécuter avec : python3 generate_presentation.py
Sortie : presentation.pptx dans le même dossier.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

# ---------------------------------------------------------------------------
# Constantes de style
# ---------------------------------------------------------------------------
W = Inches(13.33)   # Largeur slide 16:9
H = Inches(7.5)     # Hauteur slide 16:9

# Palette sobre : blanc / noir / bleu acier / gris moyen
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
C_ACCENT  = RGBColor(0x1F, 0x38, 0x64)   # Bleu foncé
C_ACCENT2 = RGBColor(0x44, 0x72, 0xC4)   # Bleu moyen
C_GRAY    = RGBColor(0x60, 0x60, 0x60)
C_LGRAY   = RGBColor(0xE8, 0xE8, 0xE8)
C_GREEN   = RGBColor(0x37, 0x86, 0x10)   # Validation
C_ORANGE  = RGBColor(0xC5, 0x50, 0x10)   # Hardware / CNC
C_RED     = RGBColor(0xC0, 0x00, 0x00)   # Alerte

FONT = "Calibri"

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width=Pt(1)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=Pt(18), bold=False, color=C_BLACK,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def set_slide_bg(slide, rgb: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_para(tf, text, font_size=Pt(16), bold=False, color=C_BLACK,
             align=PP_ALIGN.LEFT, italic=False, space_before=Pt(4), bullet=False):
    from pptx.oxml.ns import qn
    from lxml import etree
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def header_bar(slide, title, subtitle=None):
    """Barre de titre en haut de chaque slide."""
    bar = add_rect(slide, 0, 0, W, Cm(2.8), fill_rgb=C_ACCENT)
    add_text_box(slide, title,
                 Cm(0.8), Cm(0.25), Cm(28), Cm(1.4),
                 font_size=Pt(24), bold=True, color=C_WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Cm(0.8), Cm(1.55), Cm(28), Cm(0.9),
                     font_size=Pt(13), color=RGBColor(0xBD, 0xD5, 0xF0), italic=True)
    # Ligne de séparation fine sous la barre
    sep = add_rect(slide, 0, Cm(2.8), W, Pt(2), fill_rgb=C_ACCENT2)


def footer(slide, page_num, total):
    add_text_box(slide,
                 f"Machine de Dépose de Pâte Thermique — Projet stage BUT3 Informatique",
                 Cm(0.5), Cm(18.8), Cm(25), Cm(0.6),
                 font_size=Pt(9), color=C_GRAY, italic=True)
    add_text_box(slide, f"{page_num} / {total}",
                 Cm(31.5), Cm(18.8), Cm(2), Cm(0.6),
                 font_size=Pt(9), color=C_GRAY, align=PP_ALIGN.RIGHT)


def content_box(slide, left, top, width, height):
    """Zone de contenu avec bord gauche coloré."""
    accent = add_rect(slide, left, top, Pt(4), height, fill_rgb=C_ACCENT2)
    return left + Pt(10), top, width - Pt(10), height


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
TOTAL_SLIDES = 9

def slide_title(prs):
    """Slide 1 — Titre."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, C_WHITE)

    # Bande bleue en haut (40% de la hauteur)
    add_rect(slide, 0, 0, W, Inches(3.2), fill_rgb=C_ACCENT)

    add_text_box(slide,
                 "Machine de Dépose de Pâte Thermique",
                 Cm(1.5), Cm(1.2), Cm(30), Cm(2.5),
                 font_size=Pt(34), bold=True, color=C_WHITE,
                 align=PP_ALIGN.LEFT)

    add_text_box(slide,
                 "Automatisation de la dépose sur coques de calculateur automobile",
                 Cm(1.5), Cm(3.8), Cm(30), Cm(1.2),
                 font_size=Pt(18), color=C_GRAY, italic=True)

    # Ligne séparatrice
    add_rect(slide, Cm(1.5), Cm(5.3), Cm(10), Pt(2), fill_rgb=C_ACCENT2)

    infos = [
        ("Cadre",      "Projet de stage — BUT Informatique 3ème année"),
        ("Période",    "Mai → Août 2026"),
        ("Jalons",     "Soutenance blanche fin juillet · Soutenance finale fin août"),
    ]
    for i, (label, value) in enumerate(infos):
        y = Cm(5.9) + i * Cm(0.95)
        add_text_box(slide, label + " :", Cm(1.5), y, Cm(4), Cm(0.8),
                     font_size=Pt(13), bold=True, color=C_ACCENT)
        add_text_box(slide, value, Cm(5.5), y, Cm(27), Cm(0.8),
                     font_size=Pt(13), color=C_BLACK)

    footer(slide, 1, TOTAL_SLIDES)


def slide_concept(prs):
    """Slide 2 — Le concept."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    header_bar(slide, "Le concept", "Pourquoi automatiser la dépose de pâte thermique ?")
    footer(slide, 2, TOTAL_SLIDES)

    # Colonne gauche — problème
    add_text_box(slide, "Contexte industriel",
                 Cm(0.8), Cm(3.4), Cm(15), Cm(0.7),
                 font_size=Pt(15), bold=True, color=C_ACCENT)
    txb = slide.shapes.add_textbox(Cm(0.8), Cm(4.1), Cm(15), Cm(8))
    tf = txb.text_frame
    tf.word_wrap = True
    items = [
        "Les coques de calculateur automobile intègrent des composants électroniques générant de la chaleur.",
        "La pâte thermique assure le transfert de chaleur entre le composant et le dissipateur.",
        "Actuellement, la dépose est manuelle : résultat variable, non reproductible.",
        "Objectif : automatiser et tracer la dépose pour garantir la qualité et la traçabilité.",
    ]
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = "→  " + item
        run.font.name = FONT
        run.font.size = Pt(13)
        run.font.color.rgb = C_BLACK

    # Séparateur vertical
    add_rect(slide, Cm(16.2), Cm(3.2), Pt(1.5), Cm(9.5), fill_rgb=C_LGRAY)

    # Colonne droite — solution
    add_text_box(slide, "La solution",
                 Cm(16.8), Cm(3.4), Cm(16), Cm(0.7),
                 font_size=Pt(15), bold=True, color=C_ACCENT)
    txb2 = slide.shapes.add_textbox(Cm(16.8), Cm(4.1), Cm(16), Cm(8))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    items2 = [
        ("Caméra RPi (CSI)", "Capture une photo de la coque avant la dépose"),
        ("Vision ArUco", "Détecte 4 marqueurs pour calibrer la perspective"),
        ("Sélection zone", "L'opérateur dessine la zone à traiter sur l'écran tactile"),
        ("Planification", "Calcul automatique de la trajectoire de dépose"),
        ("Machine CNC", "Dépose la pâte selon la trajectoire calculée en G-code"),
        ("Rapport PDF", "Photo avant/après + paramètres, horodaté et archivé"),
    ]
    first2 = True
    for label, desc in items2:
        p = tf2.paragraphs[0] if first2 else tf2.add_paragraph()
        first2 = False
        p.space_before = Pt(6)
        r1 = p.add_run()
        r1.text = label + " : "
        r1.font.name = FONT
        r1.font.size = Pt(13)
        r1.font.bold = True
        r1.font.color.rgb = C_ACCENT2
        r2 = p.add_run()
        r2.text = desc
        r2.font.name = FONT
        r2.font.size = Pt(13)
        r2.font.color.rgb = C_BLACK


def slide_workflow(prs):
    """Slide 3 — Processus de dépose (workflow opérateur)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    header_bar(slide, "Processus de dépose", "Cycle complet — du déclenchement au rapport")
    footer(slide, 3, TOTAL_SLIDES)

    steps = [
        ("1", "Photo\nde la pièce",   "Caméra RPi capture\nla coque"),
        ("2", "Sélection\nde la zone", "L'opérateur dessine\nla zone sur l'écran"),
        ("3", "Calcul\ntrajectoire",   "PathPlanner génère\nle parcours G-code"),
        ("4", "Dépose\nautomatique",   "La machine suit\nla trajectoire"),
        ("5", "Photo\naprès dépose",  "Vérification visuelle\ndu résultat"),
        ("6", "Rapport\nPDF",          "Archivage horodaté\navant/après"),
    ]

    step_w = Cm(5.0)
    step_h = Cm(4.5)
    gap    = Cm(0.5)
    start_x = Cm(0.5)
    y_box  = Cm(3.4)

    for i, (num, title, desc) in enumerate(steps):
        x = start_x + i * (step_w + gap)

        # Boîte numéro
        num_box = add_rect(slide, x, y_box, step_w, Cm(1.0), fill_rgb=C_ACCENT)
        add_text_box(slide, num, x, y_box, step_w, Cm(1.0),
                     font_size=Pt(18), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Boîte titre
        add_rect(slide, x, y_box + Cm(1.0), step_w, Cm(1.4),
                 fill_rgb=C_ACCENT2)
        add_text_box(slide, title, x, y_box + Cm(1.0), step_w, Cm(1.4),
                     font_size=Pt(12), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Boîte description
        add_rect(slide, x, y_box + Cm(2.4), step_w, Cm(2.1),
                 fill_rgb=None, line_rgb=C_LGRAY, line_width=Pt(1))
        add_text_box(slide, desc, x + Cm(0.2), y_box + Cm(2.5), step_w - Cm(0.4), Cm(2.0),
                     font_size=Pt(11), color=C_GRAY, align=PP_ALIGN.CENTER)

        # Flèche entre les boîtes (sauf après la dernière)
        if i < len(steps) - 1:
            arrow_x = x + step_w + gap * 0.1
            add_text_box(slide, "→", arrow_x, y_box + Cm(1.3), gap + Cm(0.4), Cm(0.8),
                         font_size=Pt(18), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # Note bas de slide
    add_text_box(slide,
                 "Ce cycle peut être relancé immédiatement pour une nouvelle pièce, sans redémarrer l'application.",
                 Cm(0.8), Cm(8.4), Cm(32), Cm(0.7),
                 font_size=Pt(11), color=C_GRAY, italic=True)


def slide_architecture_hw(prs):
    """Slide 4 — Architecture matérielle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    header_bar(slide, "Architecture matérielle", "Raspberry Pi 3B+ + machine CNC (firmware Marlin)")
    footer(slide, 4, TOTAL_SLIDES)

    # Bloc Raspberry Pi
    add_rect(slide, Cm(0.8), Cm(3.2), Cm(14), Cm(7.8),
             fill_rgb=None, line_rgb=C_ACCENT, line_width=Pt(2))
    add_text_box(slide, "Raspberry Pi 3B+", Cm(1.0), Cm(3.2), Cm(13.5), Cm(0.8),
                 font_size=Pt(13), bold=True, color=C_ACCENT)
    modules = [
        "Logiciel Python (modules + GUI)",
        "Caméra CSI — module RPi (nappe 15 br.)",
        "Écran tactile 7\" 800×480 (HDMI + USB touch)",
        "Traitement ArUco (OpenCV)",
        "Interface utilisateur (PyQt5)",
    ]
    txb = slide.shapes.add_textbox(Cm(1.3), Cm(4.1), Cm(13), Cm(6.5))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for m in modules:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = "•  " + m
        run.font.name = FONT
        run.font.size = Pt(13)
        run.font.color.rgb = C_BLACK

    # Flèche centrale
    add_text_box(slide, "USB série\nG-code Marlin\n115 200 baud",
                 Cm(15.2), Cm(5.8), Cm(3.4), Cm(2.5),
                 font_size=Pt(11), color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)
    add_text_box(slide, "⟺", Cm(15.5), Cm(7.0), Cm(2.8), Cm(1.2),
                 font_size=Pt(28), bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

    # Bloc Machine CNC
    add_rect(slide, Cm(19.2), Cm(3.2), Cm(13.8), Cm(7.8),
             fill_rgb=None, line_rgb=C_ORANGE, line_width=Pt(2))
    add_text_box(slide, "Machine CNC (firmware Marlin)", Cm(19.5), Cm(3.2), Cm(13.3), Cm(0.8),
                 font_size=Pt(13), bold=True, color=C_ORANGE)
    cnc_items = [
        "Carte contrôleur — protocole G-code standard",
        "Axes X / Y — déplacement de la buse",
        "Axe Z — hauteur de la buse",
        "Axe E → piston Nema 17 — dépose pâte",
        "Fins de course — homing automatique",
        "Port série : /dev/ttyUSB0 ou /dev/ttyACM0",
    ]
    txb2 = slide.shapes.add_textbox(Cm(19.8), Cm(4.1), Cm(12.8), Cm(6.5))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    first2 = True
    for item in cnc_items:
        p = tf2.paragraphs[0] if first2 else tf2.add_paragraph()
        first2 = False
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = "•  " + item
        run.font.name = FONT
        run.font.size = Pt(13)
        run.font.color.rgb = C_BLACK

    # Marqueurs ArUco en bas
    add_rect(slide, Cm(0.8), Cm(11.3), Cm(32.2), Cm(0.9),
             fill_rgb=C_LGRAY)
    add_text_box(slide,
                 "Référentiel géométrique : 4 marqueurs ArUco (DICT_4X4_50, IDs 0–3) — calibrage de perspective par vision",
                 Cm(1.2), Cm(11.3), Cm(31.5), Cm(0.9),
                 font_size=Pt(11), color=C_GRAY, italic=True)


def slide_two_machines(prs):
    """Slide 5 — Stratégie deux machines."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    header_bar(slide, "Stratégie matérielle", "Deux machines, un seul logiciel")
    footer(slide, 5, TOTAL_SLIDES)

    # Machine PoC
    add_rect(slide, Cm(0.8), Cm(3.2), Cm(15.5), Cm(8.6),
             fill_rgb=None, line_rgb=C_ACCENT2, line_width=Pt(2))
    add_rect(slide, Cm(0.8), Cm(3.2), Cm(15.5), Cm(1.0), fill_rgb=C_ACCENT2)
    add_text_box(slide, "Machine PoC — Geeetech I3 (imprimante 3D modifiée)",
                 Cm(1.0), Cm(3.2), Cm(15.1), Cm(1.0),
                 font_size=Pt(12), bold=True, color=C_WHITE)

    poc_items = [
        ("Rôle", "Développement et validation du logiciel"),
        ("Disponibilité", "Immédiatement — déjà en place"),
        ("Firmware", "Marlin (version à identifier via M115)"),
        ("Axes", "X/Y/Z + axe E (extrudeur → piston)"),
        ("Contrôleur", "Carte Geeetech d'origine"),
        ("Avantage", "Accès immédiat, peut fonctionner en mode test"),
        ("Limite", "Pas la machine finale — sert uniquement de banc de test"),
    ]
    txb = slide.shapes.add_textbox(Cm(1.2), Cm(4.4), Cm(14.8), Cm(7.0))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for label, val in poc_items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        r1 = p.add_run()
        r1.text = label + " : "
        r1.font.name = FONT
        r1.font.size = Pt(12)
        r1.font.bold = True
        r1.font.color.rgb = C_ACCENT2
        r2 = p.add_run()
        r2.text = val
        r2.font.name = FONT
        r2.font.size = Pt(12)
        r2.font.color.rgb = C_BLACK

    # Flèche de portage
    add_text_box(slide, "Portage\nlogiciel\n→", Cm(16.5), Cm(6.2), Cm(2.0), Cm(2.4),
                 font_size=Pt(14), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, "config.py\nuniquement",
                 Cm(16.3), Cm(8.2), Cm(2.5), Cm(1.0),
                 font_size=Pt(10), color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Machine cible
    add_rect(slide, Cm(19.0), Cm(3.2), Cm(14.0), Cm(8.6),
             fill_rgb=None, line_rgb=C_ORANGE, line_width=Pt(2))
    add_rect(slide, Cm(19.0), Cm(3.2), Cm(14.0), Cm(1.0), fill_rgb=C_ORANGE)
    add_text_box(slide, "Machine cible — CNC (carte Marlin)",
                 Cm(19.2), Cm(3.2), Cm(13.6), Cm(1.0),
                 font_size=Pt(12), bold=True, color=C_WHITE)

    cnc_items = [
        ("Rôle", "Machine de production finale"),
        ("Disponibilité", "Assemblage prévu en juillet 2026"),
        ("Firmware", "Marlin (même protocole G-code)"),
        ("Axes", "X/Y/Z + actionneur de dépose"),
        ("Contrôleur", "Carte CNC dédiée"),
        ("Avantage", "Conçue pour la tâche, plus précise et robuste"),
        ("RPi + caméra", "Identiques — transférés depuis la Geeetech"),
    ]
    txb2 = slide.shapes.add_textbox(Cm(19.4), Cm(4.4), Cm(13.2), Cm(7.0))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    first2 = True
    for label, val in cnc_items:
        p = tf2.paragraphs[0] if first2 else tf2.add_paragraph()
        first2 = False
        p.space_before = Pt(5)
        r1 = p.add_run()
        r1.text = label + " : "
        r1.font.name = FONT
        r1.font.size = Pt(12)
        r1.font.bold = True
        r1.font.color.rgb = C_ORANGE
        r2 = p.add_run()
        r2.text = val
        r2.font.name = FONT
        r2.font.size = Pt(12)
        r2.font.color.rgb = C_BLACK

    # Note bas
    add_text_box(slide,
                 "Le logiciel est identique sur les deux machines — seuls le port série et les dimensions de la zone de travail changent dans config.py.",
                 Cm(0.8), Cm(12.1), Cm(32.2), Cm(0.7),
                 font_size=Pt(11), color=C_GRAY, italic=True)


def slide_architecture_sw(prs):
    """Slide 6 — Architecture logicielle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    header_bar(slide, "Architecture logicielle", "7 modules Python + interface PyQt5")
    footer(slide, 6, TOTAL_SLIDES)

    modules = [
        ("config.py",       "Paramètres globaux (ports, résolutions, constantes)",  C_ACCENT),
        ("camera.py",       "Capture image via caméra RPi (interface CSI/V4L2)",    C_ACCENT2),
        ("vision.py",       "Détection ArUco, calcul homographie, redressement",    C_ACCENT2),
        ("machine.py",      "Communication G-code Marlin via USB série",             C_ACCENT2),
        ("path_planner.py", "Calcul trajectoires de dépose (hachures, spirale…)",   C_ACCENT2),
        ("reporter.py",     "Génération rapport PDF horodaté (fpdf2)",               C_ACCENT2),
        ("main.py",         "Machine à états : IDLE→CAPTURE→CONFIGURE→RUNNING→DONE",C_ACCENT),
    ]

    col_w = Cm(15.8)
    for i, (name, desc, color) in enumerate(modules):
        row = i % 4
        col = i // 4
        x = Cm(0.8) + col * (col_w + Cm(0.6))
        y = Cm(3.4) + row * Cm(2.4)

        add_rect(slide, x, y, col_w, Cm(2.1),
                 fill_rgb=None, line_rgb=C_LGRAY, line_width=Pt(1))
        add_rect(slide, x, y, Cm(0.35), Cm(2.1), fill_rgb=color)

        add_text_box(slide, name, x + Cm(0.55), y + Cm(0.15), col_w - Cm(0.7), Cm(0.7),
                     font_size=Pt(13), bold=True, color=color)
        add_text_box(slide, desc, x + Cm(0.55), y + Cm(0.8), col_w - Cm(0.7), Cm(1.1),
                     font_size=Pt(11), color=C_GRAY)

    # GUI séparément
    add_text_box(slide, "Interface graphique (gui/)",
                 Cm(0.8), Cm(13.6), Cm(8), Cm(0.7),
                 font_size=Pt(13), bold=True, color=C_ACCENT)
    gui_items = "4 écrans PyQt5 · 800×480 · plein écran · tactile  →  Capture · Zone/Quantité · Exécution · Rapport"
    add_text_box(slide, gui_items,
                 Cm(0.8), Cm(14.2), Cm(32), Cm(0.7),
                 font_size=Pt(12), color=C_GRAY)


def slide_planning(prs):
    """Slide 7 — Planning global."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    header_bar(slide, "Planning global", "3 parties · 3 jalons · mai → août 2026")
    footer(slide, 7, TOTAL_SLIDES)

    parties = [
        {
            "titre": "Partie A — Logiciel sur Geeetech (PoC)",
            "periode": "27 mai → ~12 juillet 2026",
            "color_rect": C_ACCENT2,
            "jalon": "Jalon A : Logiciel validé sur Geeetech",
            "phases": [
                "Ph. 0  Identification matériel             1 session",
                "Ph. 1  Caméra de base                     1 session",
                "Ph. 2  ArUco & calibrage                  3 sessions",
                "Ph. 3  Communication G-code Marlin         2 sessions",
                "Ph. 4  Interface graphique                 3 sessions",
                "Ph. 5  Sélection zone & trajectoire        3 sessions",
                "Ph. 6  Intégration workflow complet        3 sessions",
                "Ph. 7  Rapport PDF                         2 sessions",
                "Ph. 8  Tests & finitions (Geeetech)        3 sessions",
            ],
            "note": "Vacances 15–19 juin · Draft rapport dû le 15 juin",
        },
        {
            "titre": "Partie B — Intégration sur CNC cible",
            "periode": "~13 juillet → fin juillet 2026",
            "color_rect": C_ORANGE,
            "jalon": "Jalon B : Soutenance blanche — fin juillet",
            "phases": [
                "Ph. 9   Assemblage CNC (hardware)          ~2 semaines",
                "Ph. 10  Portage logiciel (config.py)       2 sessions",
                "Ph. 11  Validation complète sur CNC        3 sessions",
            ],
            "note": "RPi, caméra et écran réutilisés sans modification",
        },
        {
            "titre": "Partie C — Finalisation",
            "periode": "Août 2026",
            "color_rect": C_GREEN,
            "jalon": "Jalon C : Soutenance finale — fin août",
            "phases": [
                "Ph. 12  Corrections bugs (retours S. blanche)  ~1 semaine",
                "Ph. 13  Finalisation rapport                   ~3 semaines",
            ],
            "note": "Rédaction rapport en parallèle depuis mai (~1h/soir)",
        },
    ]

    col_w = Cm(10.5)
    col_gap = Cm(0.45)
    x_start = Cm(0.6)

    for i, p in enumerate(parties):
        x = x_start + i * (col_w + col_gap)
        y = Cm(3.1)

        # En-tête de colonne
        add_rect(slide, x, y, col_w, Cm(1.0), fill_rgb=p["color_rect"])
        add_text_box(slide, p["titre"], x + Cm(0.2), y, col_w - Cm(0.4), Cm(1.0),
                     font_size=Pt(11), bold=True, color=C_WHITE)

        add_text_box(slide, p["periode"],
                     x, y + Cm(1.05), col_w, Cm(0.55),
                     font_size=Pt(10), color=C_GRAY, italic=True)

        # Liste des phases
        txb = slide.shapes.add_textbox(x, y + Cm(1.65), col_w, Cm(7.5))
        tf = txb.text_frame
        tf.word_wrap = True
        first = True
        for phase in p["phases"]:
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            para.space_before = Pt(4)
            run = para.add_run()
            run.text = "•  " + phase
            run.font.name = FONT
            run.font.size = Pt(10)
            run.font.color.rgb = C_BLACK

        # Note bas de colonne
        add_text_box(slide, p["note"],
                     x, y + Cm(9.4), col_w, Cm(0.8),
                     font_size=Pt(10), color=C_GRAY, italic=True)

        # Jalon
        add_rect(slide, x, y + Cm(10.3), col_w, Cm(0.75), fill_rgb=RGBColor(0xFF, 0xD7, 0x00))
        add_text_box(slide, "★  " + p["jalon"],
                     x + Cm(0.2), y + Cm(10.3), col_w - Cm(0.3), Cm(0.75),
                     font_size=Pt(10), bold=True, color=RGBColor(0x5A, 0x37, 0x00))


def slide_jalons(prs):
    """Slide 8 — Jalons clés."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    header_bar(slide, "Jalons clés du projet", "4 dates structurantes")
    footer(slide, 8, TOTAL_SLIDES)

    jalons = [
        {
            "id": "JD",
            "label": "Premier draft du rapport",
            "date": "15 juin 2026",
            "color": C_RED,
            "critere": "Document complet (même incomplet sur certains points) remis avant les vacances.",
            "note": "⚠️ Tombe la même semaine que la Phase 4 (GUI) — gérer le temps semaine du 8 au 14 juin.",
        },
        {
            "id": "J1",
            "label": "Logiciel validé sur Geeetech",
            "date": "~12 juillet 2026",
            "color": C_ACCENT2,
            "critere": "Cycle complet (photo → sélection → dépose → PDF) fonctionnel sur la Geeetech.",
            "note": "Suite : transfert du Raspberry Pi et de la caméra vers la CNC cible.",
        },
        {
            "id": "J2",
            "label": "Soutenance blanche",
            "date": "Fin juillet 2026",
            "color": C_ORANGE,
            "critere": "Démo live sur CNC cible. Rapport en version avancée.",
            "note": "Retours à intégrer en août (corrections Phase 12).",
        },
        {
            "id": "J3",
            "label": "Soutenance finale",
            "date": "Fin août 2026",
            "color": C_GREEN,
            "critere": "Rapport remis. Présentation et démonstration devant le jury.",
            "note": "Fin du projet.",
        },
    ]

    box_h = Cm(2.7)
    box_w = Cm(31.5)
    x0 = Cm(1.0)
    y0 = Cm(3.3)

    for i, j in enumerate(jalons):
        y = y0 + i * (box_h + Cm(0.35))

        # Pastille ID
        add_rect(slide, x0, y, Cm(2.0), box_h, fill_rgb=j["color"])
        add_text_box(slide, j["id"], x0, y, Cm(2.0), box_h,
                     font_size=Pt(16), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Contenu
        add_rect(slide, x0 + Cm(2.0), y, box_w - Cm(2.0), box_h,
                 fill_rgb=None, line_rgb=C_LGRAY, line_width=Pt(1))
        add_text_box(slide, j["label"] + "  —  " + j["date"],
                     x0 + Cm(2.3), y + Cm(0.1), box_w - Cm(2.5), Cm(0.75),
                     font_size=Pt(14), bold=True, color=C_BLACK)
        add_text_box(slide, "Critère : " + j["critere"],
                     x0 + Cm(2.3), y + Cm(0.85), box_w - Cm(2.5), Cm(0.75),
                     font_size=Pt(11), color=C_GRAY)
        add_text_box(slide, j["note"],
                     x0 + Cm(2.3), y + Cm(1.6), box_w - Cm(2.5), Cm(0.75),
                     font_size=Pt(11), color=j["color"], italic=True)


def slide_conclusion(prs):
    """Slide 9 — Synthèse."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    header_bar(slide, "Synthèse", "Ce qu'il faut retenir")
    footer(slide, 9, TOTAL_SLIDES)

    points = [
        (C_ACCENT2,  "Concept",
         "Automatiser la dépose de pâte thermique sur coques de calculateur automobile "
         "avec vision par ordinateur (ArUco), machine CNC (Marlin) et interface tactile (PyQt5)."),
        (C_ACCENT2,  "Approche en deux machines",
         "La Geeetech I3 (PoC) sert à développer et valider le logiciel. "
         "La CNC cible est la machine de production finale. "
         "Le portage est minimal : seul config.py change."),
        (C_ORANGE,   "Planning resserré",
         "Développement logiciel de fin mai à mi-juillet. "
         "Assemblage CNC + validation : juillet. "
         "Rapport + soutenance : août."),
        (C_RED,      "Point d'attention",
         "Semaine du 8–14 juin : Phase 4 (GUI) ET deadline draft rapport (15 juin) en même temps. "
         "Anticiper la rédaction du rapport dès la semaine du 1er juin."),
        (C_GREEN,    "Clé de réussite",
         "Maintenir CONCEPTION.md à jour après chaque session : "
         "il constitue le brouillon permanent du rapport de soutenance."),
    ]

    y0 = Cm(3.3)
    for i, (color, title, body) in enumerate(points):
        y = y0 + i * Cm(2.2)
        add_rect(slide, Cm(0.8), y, Cm(0.4), Cm(1.8), fill_rgb=color)
        add_text_box(slide, title, Cm(1.5), y, Cm(5.5), Cm(0.75),
                     font_size=Pt(13), bold=True, color=color)
        add_text_box(slide, body, Cm(1.5), y + Cm(0.75), Cm(31.2), Cm(1.2),
                     font_size=Pt(12), color=C_BLACK)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_concept(prs)
    slide_workflow(prs)
    slide_architecture_hw(prs)
    slide_two_machines(prs)
    slide_architecture_sw(prs)
    slide_planning(prs)
    slide_jalons(prs)
    slide_conclusion(prs)

    output = "presentation.pptx"
    prs.save(output)
    print(f"Présentation générée : {output}  ({TOTAL_SLIDES} slides)")

if __name__ == "__main__":
    main()
